"""
Builds a concise presentation deck (~15 slides) for the Multi-Agent Research
System, drawn from the internship / mini-project report. 16:9, clean light
theme with a navy/orange accent, reusing the figure images in assets/.

    python build_ppt.py  ->  Multi_Agent_Research_System_PPT_Suraj_Singh_Rawat.pptx
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(ROOT, "assets")
OUTPUT = os.path.join(ROOT, "Multi_Agent_Research_System_PPT_Suraj_Singh_Rawat.pptx")

STUDENT = "Suraj Singh Rawat"
ROLL = "2400951530054"
CLASS = "TT-C (AI & ML)"
TITLE = "Multi-Agent Research System"
SUBTITLE = "A Multi-Agent AI Orchestration System for Automated Desk Research"

# palette
NAVY = RGBColor(0x1F, 0x2D, 0x50)
BLUE = RGBColor(0x2E, 0x52, 0x8A)
ORANGE = RGBColor(0xE0, 0x6A, 0x1A)
INK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x60, 0x60, 0x60)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW = prs.slide_width
SH = prs.slide_height


def _set(frame, text, size, *, bold=False, color=INK, align=PP_ALIGN.LEFT,
         font=FONT, space_after=6):
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = align
    para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return para


def textbox(slide, x, y, w, h, text, size, **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set(box.text_frame, text, size, **kw)
    return box


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def header(slide, kicker, title):
    accent_bar(slide)
    textbox(slide, 0.7, 0.45, 11.5, 0.4, kicker.upper(), 13,
            bold=True, color=ORANGE)
    textbox(slide, 0.68, 0.85, 12.0, 1.0, title, 32, bold=True, color=NAVY)
    ln = slide.shapes.add_shape(1, Inches(0.75), Inches(1.75), Inches(11.8),
                                Pt(2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor(0xD0, 0xD7, 0xE2)
    ln.line.fill.background()


def bullets(slide, items, *, x=0.9, y=2.1, w=11.6, h=4.9, size=18, gap=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        level = 0
        if isinstance(item, tuple):
            item, level = item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(gap if level == 0 else gap - 4)
        para.level = level
        run = para.add_run()
        run.text = ("•  " if level == 0 else "–  ") + item
        run.font.size = Pt(size if level == 0 else size - 3)
        run.font.name = FONT
        run.font.color.rgb = INK if level == 0 else GREY
        run.font.bold = False
    return box


def picture(slide, filename, *, x, y, w=None, h=None):
    path = os.path.join(ASSETS, filename)
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)


def table_slide(slide, headers, rows, *, x=0.9, y=2.1, w=11.6, col_widths=None,
                font_size=14):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = Inches(0.5 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                   Inches(w), height)
    tbl = shape.table
    if col_widths:
        for idx, cw in enumerate(col_widths):
            tbl.columns[idx].width = Inches(cw)
    for c, text in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(font_size + 1)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(text)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.name = FONT
            p.runs[0].font.color.rgb = INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl


def footer(slide, n):
    textbox(slide, 0.7, 7.02, 8.0, 0.35,
            f"{TITLE}  ·  {STUDENT}", 9, color=GREY)
    textbox(slide, 12.3, 7.02, 0.8, 0.35, str(n), 9, color=GREY,
            align=PP_ALIGN.RIGHT)


def content_slide(kicker, title, n):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, WHITE)
    header(slide, kicker, title)
    footer(slide, n)
    return slide


# ----------------------------------------------------------------------------
# 1  Title
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
band = s.shapes.add_shape(1, Inches(0), Inches(4.55), SW, Inches(0.08))
band.fill.solid()
band.fill.fore_color.rgb = ORANGE
band.line.fill.background()

textbox(s, 0.9, 1.5, 11.5, 0.5, "MINI PROJECT PRESENTATION", 16, bold=True,
        color=RGBColor(0xF2, 0xA6, 0x5E))
textbox(s, 0.85, 2.1, 11.7, 1.6, TITLE, 52, bold=True, color=WHITE)
textbox(s, 0.9, 3.75, 11.5, 0.8, SUBTITLE, 20,
        color=RGBColor(0xC8, 0xD2, 0xE6))

textbox(s, 0.9, 4.95, 11.5, 0.45, f"Submitted by:  {STUDENT}", 18, bold=True,
        color=WHITE)
textbox(s, 0.9, 5.45, 11.5, 0.4, f"University Roll No: {ROLL}   |   Class: {CLASS}",
        14, color=RGBColor(0xC8, 0xD2, 0xE6))
textbox(s, 0.9, 6.05, 11.5, 0.8,
        "Department of Computer Science & Engineering\n"
        "MGM's College of Engineering & Technology, Noida  ·  2026-27", 12,
        color=RGBColor(0xAE, 0xBB, 0xD4))
try:
    s.shapes.add_picture(os.path.join(ASSETS, "mgm_logo.jpeg"),
                         Inches(11.3), Inches(0.55), height=Inches(1.15))
except Exception:
    pass

# ----------------------------------------------------------------------------
# 2  Problem
# ----------------------------------------------------------------------------
s = content_slide("Motivation", "The Problem: Desk Research Is Slow and Manual", 2)
bullets(s, [
    "Preparing a briefing on an unfamiliar topic follows a fixed sequence: "
    "find credible sources, read them, condense them into prose, judge if it is complete.",
    "That sequence takes 2-4 hours per topic for a human analyst.",
    "Output quality varies widely with the individual analyst's skill and diligence.",
    "Capacity scales only by hiring more people — cost rises in direct proportion to volume.",
    ("Because the sequence is predictable, it is a strong candidate for automation.", 1),
])

# ----------------------------------------------------------------------------
# 3  Objectives
# ----------------------------------------------------------------------------
s = content_slide("Objectives", "Design Goals — Agreed and Numeric", 3)
bullets(s, [
    "Perform the whole desk-research sequence with no human intervention.",
    "Speed  —  produce a complete report within 60 seconds.",
    "Cost  —  substantially below the commercial LLM endpoint then in use.",
    "Scale  —  serve at least 10 simultaneous users without degradation.",
    "Gain practical experience of the engineering around ML work: automated "
    "testing, benchmarking and deployment — not model development in isolation.",
])
textbox(s, 0.9, 6.15, 11.6, 0.6,
        "Success is demonstrated by measurement, not asserted in prose.", 15,
        bold=True, color=ORANGE)

# ----------------------------------------------------------------------------
# 4  Techniques & tools
# ----------------------------------------------------------------------------
s = content_slide("Background", "Techniques and Tools Studied", 4)
bullets(s, [
    "Large Language Models  —  stateless; everything the model must consider "
    "has to be in the prompt. Instruction phrasing that fixes output shape beats "
    "phrasing that asks for quality.",
    "The Multi-Agent Pattern  —  split a task into single-purpose stages with "
    "declared inputs/outputs, instead of one opaque prompt.",
    "LangChain  —  templated prompts, provider abstraction, output parsing, "
    "retry behaviour.",
    "Web retrieval & extraction  —  Tavily / DuckDuckGo for discovery, "
    "BeautifulSoup4 for stripping markup down to body text.",
    "Streamlit  —  renders a Python script as a web app; simple UI, no second "
    "language or build step.",
])

# ----------------------------------------------------------------------------
# 5  Architecture
# ----------------------------------------------------------------------------
s = content_slide("Design", "System Architecture — Four Layers", 5)
picture(s, "fig_3_1_architecture.png", x=1.55, y=2.05, h=4.5)
textbox(s, 0.9, 6.7, 11.6, 0.5,
        "Presentation  →  Orchestration  →  Agents  →  Services.   "
        "No layer reaches past its neighbour, so a change in one (e.g. swapping the "
        "model provider) does not propagate.", 12, color=GREY)

# ----------------------------------------------------------------------------
# 6  Pipeline
# ----------------------------------------------------------------------------
s = content_slide("Design", "The Four-Agent Sequential Pipeline", 6)
picture(s, "fig_3_2_pipeline.png", x=2.4, y=2.0, h=4.6)
textbox(s, 0.9, 6.75, 11.6, 0.5,
        "Fixed order: reading needs sources, writing needs material, evaluation "
        "needs a report. Measured mean end-to-end: 49.1 s against a 60 s ceiling.",
        12, color=GREY)

# ----------------------------------------------------------------------------
# 7  Agent roles
# ----------------------------------------------------------------------------
s = content_slide("Design", "What Each Agent Does", 7)
table_slide(s,
            ["Agent", "Receives", "Returns"],
            [["Search", "Topic string",
              "Up to 10 sources, each with URL, title, snippet, relevance score"],
             ["Reader", "Ranked sources",
              "Aggregated body text + counts of successful / failed retrievals"],
             ["Writer", "Aggregated text + topic",
              "Report in Markdown, six fixed sections"],
             ["Critic", "Report + topic",
              "Score 0-10, five sub-scores, written strengths & weaknesses"]],
            col_widths=[1.7, 3.2, 6.7], font_size=13)
textbox(s, 0.9, 5.0, 11.6, 0.6,
        "Each agent has a declared contract — so each can be timed, tested and "
        "fixed on its own, and a defect localises to one stage.", 14,
        bold=True, color=BLUE)

# ----------------------------------------------------------------------------
# 8  Tech stack
# ----------------------------------------------------------------------------
s = content_slide("Design", "Technology Stack and Rationale", 8)
table_slide(s,
            ["Layer", "Technology", "Reason for selection"],
            [["Interface", "Streamlit", "Simple UI; no second language or build step"],
             ["Orchestration", "LangChain", "Prompt templating, output parsing, provider abstraction"],
             ["Language model", "Groq gpt-oss-120b", "~70% cheaper per token; very low latency"],
             ["Discovery", "Tavily / DuckDuckGo", "No quota at the volumes involved"],
             ["Extraction", "BeautifulSoup4", "Tolerant of malformed markup; well documented"],
             ["Config", "python-dotenv", "Keeps credentials out of the repository"],
             ["Testing", "pytest", "Fixtures and parametrisation suit agent testing"]],
            col_widths=[1.9, 2.7, 7.0], font_size=13)

# ----------------------------------------------------------------------------
# 9  Implementation highlights
# ----------------------------------------------------------------------------
s = content_slide("Implementation", "Implementation Highlights", 9)
bullets(s, [
    "Prompt design  —  the Writer prompt names six sections with word budgets; "
    "stating structure explicitly was the single biggest gain in output consistency.",
    "Sampling  —  temperature 0.7 for synthesis (readable prose), 0.3 for the "
    "Critic (consistent scoring).",
    "Reader is built to fail gracefully  —  a dead URL removes one source, it "
    "does not end the request; extractions under 100 words are rejected.",
    "Error handling  —  page failure absorbed by Reader; model-endpoint failure "
    "retried with exponential backoff (2 s → 32 s); missing credential is terminal.",
    "Stateless  —  each request owns its state; no database; any instance can be "
    "restarted without data loss.",
])

# ----------------------------------------------------------------------------
# 10  UI
# ----------------------------------------------------------------------------
s = content_slide("Implementation", "User Interface", 10)
picture(s, "fig_4_3_interface.png", x=3.2, y=2.0, h=4.5)
textbox(s, 0.9, 6.65, 11.6, 0.6,
        "One text field, a button, a live progress indicator, the rendered report "
        "with its quality score and every source listed. All generated content is "
        "escaped before rendering.", 12, color=GREY)

# ----------------------------------------------------------------------------
# 11  Testing & performance
# ----------------------------------------------------------------------------
s = content_slide("Results", "Testing and Performance", 11)
picture(s, "fig_5_1_performance.png", x=7.0, y=2.1, h=3.9)
bullets(s, [
    "39 automated tests — unit, integration, end-to-end — all pass.",
    "Statement coverage: 85% overall (target met).",
    "Every stage finishes inside its budget.",
    "Pipeline mean: 49.1 s vs 60 s ceiling; 95th pct 58.1 s.",
    "10 concurrent users add only ~1.4 s of latency.",
], x=0.9, y=2.2, w=6.0, size=16, gap=12)

# ----------------------------------------------------------------------------
# 12  Cost
# ----------------------------------------------------------------------------
s = content_slide("Results", "Cost Evaluation", 12)
table_slide(s,
            ["Item", "Commercial endpoint", "Groq endpoint"],
            [["Model charges / month", "5.25 USD", "1.75 USD"],
             ["Compute instance", "30.00 USD", "30.00 USD"],
             ["Storage and transfer", "2.00 USD", "2.00 USD"],
             ["Monthly total", "37.25 USD", "33.75 USD"],
             ["Cost per report", "0.0373 USD", "0.0338 USD"]],
            col_widths=[3.6, 4.0, 4.0], font_size=14)
textbox(s, 0.9, 5.4, 11.6, 1.0,
        "Model charges fell ~70% (the target). Effect on the total is smaller "
        "because the compute instance dominates at 1,000 reports/month — the "
        "saving grows as volume rises, since model cost scales with usage and the "
        "instance cost does not.", 14, color=INK)

# ----------------------------------------------------------------------------
# 13  Limitations
# ----------------------------------------------------------------------------
s = content_slide("Evaluation", "Known Limitations", 13)
bullets(s, [
    "No fact-checking  —  the system synthesises what its sources assert; wrong "
    "sources produce a confident, well-formed, wrong report.",
    "The Critic scores craftsmanship, not correctness  —  it has no ground truth.",
    "Extraction is heuristic  —  fails on script-rendered pages, paywalls and "
    "non-HTML; typically 1-3 of 10 URLs are lost.",
    "Pipeline is strictly sequential  —  it cannot revisit an earlier stage in "
    "light of a later one.",
    "English only; nothing stored, so no report history; Critic Agent coverage is "
    "79% vs the 85% target.",
])

# ----------------------------------------------------------------------------
# 14  Conclusion & future scope
# ----------------------------------------------------------------------------
s = content_slide("Conclusion", "Conclusion and Future Scope", 14)
bullets(s, [
    "The system performs the whole desk-research sequence unattended and meets "
    "every target: 49.1 s mean, ~70% lower model cost, 10 concurrent users, 39/39 tests.",
    "The decisive decision was four agents over one model — the gain was "
    "diagnosability, not just quality. Nearly every measurement depends on it.",
    ("Future scope:", 0),
    ("Act on the Critic's score — low score triggers a fresh search (line → loop).", 1),
    ("Better extraction via a headless browser or article-extraction library.", 1),
    ("Cache results by topic; add multi-language and PDF support.", 1),
])

# ----------------------------------------------------------------------------
# 15  Thank you
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
band = s.shapes.add_shape(1, Inches(0), Inches(4.35), SW, Inches(0.08))
band.fill.solid()
band.fill.fore_color.rgb = ORANGE
band.line.fill.background()
textbox(s, 0.9, 2.7, 11.5, 1.2, "Thank You", 54, bold=True, color=WHITE)
textbox(s, 0.9, 4.6, 11.5, 0.5, "Questions & Discussion", 22,
        color=RGBColor(0xC8, 0xD2, 0xE6))
textbox(s, 0.9, 5.6, 11.5, 0.5, f"{STUDENT}  ·  Roll No. {ROLL}  ·  {CLASS}",
        14, color=RGBColor(0xAE, 0xBB, 0xD4))


prs.save(OUTPUT)
print("wrote", OUTPUT)
